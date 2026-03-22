from pptx.util import Emu


def clean_text(text):
    """Remove junk characters and extra whitespace."""
    import re
    text = text.strip()
    text = re.sub(r'[^\x20-\x7E\n]', '', text)   # remove non-printable chars
    text = re.sub(r'[ \t]+', ' ', text)          # collapse multiple spaces
    text = re.sub(r'\n+', '\n', text)            # collapse multiple newlines
    return text.strip()


def emu_to_px(emu_value):
    """Convert EMU (PowerPoint units) to pixels at 96 DPI."""
    return round(emu_value / 914400 * 96)


def extract_slide_text(prs):
    """
    Existing function (DO NOT MODIFY - used by Member 4)

    Returns:
    {
      1: [
            {
              "text": "...",
              "type": "...",
              "x": ..., "y": ..., "w": ..., "h": ...
            }
         ]
    }
    """
    slide_data = {}

    for i, slide in enumerate(prs.slides):
        slide_number = i + 1
        elements = []

        for shape in slide.shapes:

            # --- Get shape coordinates ---
            x = emu_to_px(shape.left)   if shape.left   else 0
            y = emu_to_px(shape.top)    if shape.top    else 0
            w = emu_to_px(shape.width)  if shape.width  else 0
            h = emu_to_px(shape.height) if shape.height else 0

            # --- Text shapes ---
            if shape.has_text_frame:
                for para_idx, para in enumerate(shape.text_frame.paragraphs):
                    line = clean_text(para.text)
                    if not line:
                        continue

                    # Detect type
                    shape_type = "title" if shape.shape_type == 13 else "bullet"
                    if hasattr(shape, "name") and "title" in shape.name.lower():
                        shape_type = "title"
                    elif para_idx == 0 and shape.text_frame.paragraphs[0].runs:
                        runs = shape.text_frame.paragraphs[0].runs
                        if runs and runs[0].font.bold:
                            shape_type = "heading"

                    elements.append({
                        "text": line,
                        "type": shape_type,
                        "x": x,
                        "y": y,
                        "w": w,
                        "h": h
                    })

            # --- Chart shapes ---
            if shape.has_chart:
                chart = shape.chart
                try:
                    for series in chart.series:
                        label = clean_text(str(series.name))
                        if label:
                            elements.append({
                                "text": label,
                                "type": "chart_label",
                                "x": x, "y": y, "w": w, "h": h
                            })
                except Exception:
                    pass

        slide_data[slide_number] = elements
        print(f"Slide {slide_number}: {len(elements)} elements extracted")

    return slide_data

def extract_context(slides_data):
    """
    Converts structured slide data (from ppt_reader.py)
    into AI-friendly context strings.

    Input:
    [
        {
            "slide_number": 1,
            "title": "...",
            "content": [...],
            "notes": "...",
            "full_text": "..."
        }
    ]

    Output:
    {
        1: "Slide 1 is about ...",
        2: "Slide 2 explains ..."
    }
    """

    context_data = {}

    for slide in slides_data:
        slide_number = slide["slide_number"]

        title = clean_text(slide["title"])
        content = [clean_text(c) for c in slide["content"] if c]
        notes = clean_text(slide["notes"]) if slide["notes"] else ""

        # Build context string
        context = f"Slide {slide_number} is about {title}. "

        if content:
            context += "It covers " + ", ".join(content) + ". "

        if notes:
            context += f"Additional explanation: {notes}."

        context_data[slide_number] = context.strip()

    return context_data