import os
import subprocess
import platform
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE


# ─────────────────────────────────────────────
# HELPERS
# ─────────────────────────────────────────────

def px_to_emu(px):
    return int(px * 914400 / 96)

def emu_to_px(emu):
    return round(emu / 914400 * 96)

def open_file(path):
    try:
        if platform.system() == "Windows":
            os.startfile(path)
        elif platform.system() == "Darwin":
            subprocess.call(["open", path])
        else:
            subprocess.call(["xdg-open", path])
        print(f"Opened: {path}")
    except Exception as e:
        print(f"Could not auto-open: {e}")
        print(f"Manually open: {path}")

def get_save_path(ppt_path):
    """
    Always saves the highlighted file to the project folder
    on the Desktop — avoids OneDrive permission errors.
    """
    project_dir = os.path.join(
        os.path.expanduser("~"),
        "OneDrive", "Desktop", "slide"
    )
    os.makedirs(project_dir, exist_ok=True)
    filename  = os.path.basename(ppt_path).replace(
        ".pptx", "_chart_highlighted.pptx"
    )
    return os.path.join(project_dir, filename)


# ─────────────────────────────────────────────
# STEP 1: Detect chart type
# ─────────────────────────────────────────────

def detect_chart_type(chart):
    """
    Returns 'bar', 'line', 'pie', 'area', 'scatter', or 'other'
    """
    try:
        ct = chart.chart_type
        type_map = {
            XL_CHART_TYPE.BAR_CLUSTERED:    "bar",
            XL_CHART_TYPE.BAR_STACKED:      "bar",
            XL_CHART_TYPE.COLUMN_CLUSTERED: "bar",
            XL_CHART_TYPE.COLUMN_STACKED:   "bar",
            XL_CHART_TYPE.LINE:             "line",
            XL_CHART_TYPE.LINE_MARKERS:     "line",
            XL_CHART_TYPE.LINE_STACKED:     "line",
            XL_CHART_TYPE.PIE:              "pie",
            XL_CHART_TYPE.PIE_EXPLODED:     "pie",
            XL_CHART_TYPE.DOUGHNUT:         "pie",
            XL_CHART_TYPE.AREA:             "area",
            XL_CHART_TYPE.XY_SCATTER:       "scatter",
        }
        return type_map.get(ct, "other")
    except Exception:
        return "other"


# ─────────────────────────────────────────────
# STEP 2: Extract chart data from one slide
# ─────────────────────────────────────────────

def extract_charts_from_slide(slide, slide_number):
    """
    Scans one slide and returns all charts as a list of dicts:
    {
        "slide_number": 3,
        "shape_name":   "Content Placeholder 5",
        "chart_type":   "bar" / "pie" / "line" / ...,
        "x", "y", "w", "h": position in pixels,
        "series": [
            {
                "name":       "Series 1",
                "values":     [1.2, 4.5, 3.0, 4.5],
                "categories": ["Point 1", "Point 2", ...]
            }
        ]
    }
    """
    charts = []

    for shape in slide.shapes:
        if not shape.has_chart:
            continue

        chart      = shape.chart
        chart_type = detect_chart_type(chart)

        x = emu_to_px(shape.left)   if shape.left   else 0
        y = emu_to_px(shape.top)    if shape.top    else 0
        w = emu_to_px(shape.width)  if shape.width  else 0
        h = emu_to_px(shape.height) if shape.height else 0

        series_data = []
        try:
            for series in chart.series:

                # Extract numeric values safely
                try:
                    values = [
                        float(v) if v is not None else 0.0
                        for v in series.values
                    ]
                except Exception:
                    values = []

                # Extract category labels safely
                try:
                    cats = chart.plots[0].series[0].data_labels
                    categories = [str(c) for c in cats]
                except Exception:
                    categories = [
                        f"Point {i+1}" for i in range(len(values))
                    ]

                series_data.append({
                    "name":       str(series.name) if series.name else "Series",
                    "values":     values,
                    "categories": categories
                })

        except Exception as e:
            print(f"  Could not read series on slide {slide_number}: {e}")

        charts.append({
            "slide_number": slide_number,
            "shape_name":   shape.name,
            "chart_type":   chart_type,
            "x": x, "y": y, "w": w, "h": h,
            "series":       series_data
        })

        print(f"Slide {slide_number}: Found {chart_type} chart — '{shape.name}'")

    return charts


# ─────────────────────────────────────────────
# STEP 3: Analyse chart — find important points
# ─────────────────────────────────────────────

def analyse_chart(chart_info):
    """
    Finds these patterns in chart data:
        peak          → highest value point
        trough        → lowest value point
        trend_up      → biggest single rise between two consecutive points
        trend_down    → biggest single drop between two consecutive points
        biggest_slice → largest slice (pie charts only)

    Returns list of finding dicts:
    {
        "type":        "peak",
        "series_name": "Series 1",
        "index":       3,
        "value":       4.5,
        "label":       "Point 4",
        "description": "Highest value: 4.5 at Point 4"
    }
    """
    findings   = []
    chart_type = chart_info.get("chart_type", "other")

    for series in chart_info.get("series", []):
        values     = series.get("values", [])
        categories = series.get("categories", [])
        name       = series.get("name", "Series")

        if not values or len(values) < 2:
            continue

        # Pad categories if shorter than values
        while len(categories) < len(values):
            categories.append(f"Point {len(categories) + 1}")

        # ── Peak ──────────────────────────────────────
        max_val = max(values)
        max_idx = values.index(max_val)
        findings.append({
            "type":        "peak",
            "series_name": name,
            "index":       max_idx,
            "value":       max_val,
            "label":       categories[max_idx],
            "description": f"Highest value: {max_val} at {categories[max_idx]}"
        })

        # ── Trough ────────────────────────────────────
        min_val = min(values)
        min_idx = values.index(min_val)
        findings.append({
            "type":        "trough",
            "series_name": name,
            "index":       min_idx,
            "value":       min_val,
            "label":       categories[min_idx],
            "description": f"Lowest value: {min_val} at {categories[min_idx]}"
        })

        # ── Biggest rise ──────────────────────────────
        biggest_rise = 0
        rise_start   = 0
        rise_end     = 1
        for i in range(len(values) - 1):
            diff = values[i + 1] - values[i]
            if diff > biggest_rise:
                biggest_rise = diff
                rise_start   = i
                rise_end     = i + 1

        if biggest_rise > 0:
            findings.append({
                "type":        "trend_up",
                "series_name": name,
                "index":       rise_end,
                "value":       biggest_rise,
                "label":       f"{categories[rise_start]} → {categories[rise_end]}",
                "description": f"Biggest rise: +{biggest_rise} "
                               f"from {categories[rise_start]} "
                               f"to {categories[rise_end]}"
            })

        # ── Biggest fall ──────────────────────────────
        biggest_fall = 0
        fall_start   = 0
        fall_end     = 1
        for i in range(len(values) - 1):
            diff = values[i] - values[i + 1]
            if diff > biggest_fall:
                biggest_fall = diff
                fall_start   = i
                fall_end     = i + 1

        if biggest_fall > 0:
            findings.append({
                "type":        "trend_down",
                "series_name": name,
                "index":       fall_end,
                "value":       biggest_fall,
                "label":       f"{categories[fall_start]} → {categories[fall_end]}",
                "description": f"Biggest drop: -{biggest_fall} "
                               f"from {categories[fall_start]} "
                               f"to {categories[fall_end]}"
            })

        # ── Biggest pie slice ─────────────────────────
        if chart_type == "pie":
            total       = sum(values) if sum(values) > 0 else 1
            biggest     = max(values)
            biggest_idx = values.index(biggest)
            pct         = round(biggest / total * 100, 1)
            findings.append({
                "type":        "biggest_slice",
                "series_name": name,
                "index":       biggest_idx,
                "value":       biggest,
                "label":       categories[biggest_idx],
                "description": f"Biggest slice: {categories[biggest_idx]} ({pct}%)"
            })

    return findings


# ─────────────────────────────────────────────
# STEP 4: Draw colored highlight on the slide
# ─────────────────────────────────────────────

def draw_highlight_on_chart(slide, chart_info, finding):
    """
    Draws a colored border around the chart and a text label.

    Color scheme:
        peak          → green
        trough        → red
        trend_up      → dark green
        trend_down    → dark red
        biggest_slice → yellow
    """
    color_map = {
        "peak":          RGBColor(0x00, 0xCC, 0x00),
        "trough":        RGBColor(0xFF, 0x33, 0x33),
        "trend_up":      RGBColor(0x00, 0x99, 0x00),
        "trend_down":    RGBColor(0xFF, 0x00, 0x00),
        "biggest_slice": RGBColor(0xFF, 0xCC, 0x00),
    }

    finding_type = finding.get("type", "peak")
    color        = color_map.get(finding_type, RGBColor(0xFF, 0x00, 0x00))

    cx = px_to_emu(chart_info["x"])
    cy = px_to_emu(chart_info["y"])
    cw = px_to_emu(chart_info["w"])
    ch = px_to_emu(chart_info["h"])

    # Colored border rectangle around the whole chart
    border            = slide.shapes.add_shape(1, cx, cy, cw, ch)
    border.fill.background()
    border.line.color.rgb = color
    border.line.width     = Pt(4)

    # Text label at top-left of chart
    label_h           = px_to_emu(36)
    label_w           = px_to_emu(min(500, chart_info["w"]))
    label_box         = slide.shapes.add_textbox(cx, cy, label_w, label_h)
    tf                = label_box.text_frame
    tf.word_wrap      = True
    p                 = tf.paragraphs[0]
    p.text            = finding.get("description", "")
    run               = p.runs[0]
    run.font.size     = Pt(11)
    run.font.bold     = True
    run.font.color.rgb = color

    print(f"  Highlight drawn: {finding['description']}")


# ─────────────────────────────────────────────
# STEP 5: Master function — scan + highlight + save + open
# ─────────────────────────────────────────────

def highlight_charts_in_ppt(ppt_path, finding_types=None):
    """
    Scans every slide, finds charts, draws highlights, saves and opens.

    finding_types:
        None                → highlight everything
        ["peak"]            → only highest values
        ["trough"]          → only lowest values
        ["trend_up"]        → only rising trends
        ["trend_down"]      → only falling trends
        ["biggest_slice"]   → only biggest pie slice
    """
    if finding_types is None:
        finding_types = [
            "peak", "trough",
            "trend_up", "trend_down",
            "biggest_slice"
        ]

    prs              = Presentation(ppt_path)
    total_highlights = 0

    for i, slide in enumerate(prs.slides):
        slide_number = i + 1
        charts       = extract_charts_from_slide(slide, slide_number)

        if not charts:
            print(f"Slide {slide_number}: No charts found")
            continue

        for chart_info in charts:
            print(f"\nAnalysing chart on Slide {slide_number}...")
            findings = analyse_chart(chart_info)

            if not findings:
                print("  No significant patterns found.")
                continue

            for finding in findings:
                if finding["type"] in finding_types:
                    draw_highlight_on_chart(slide, chart_info, finding)
                    total_highlights += 1

    if total_highlights == 0:
        print("\nNo charts with data found in this PPT.")
        print("Make sure charts are real PowerPoint objects")
        print("(Insert → Chart), not screenshots or images.")
        return None

    # ── Save to project folder, NOT OneDrive Documents ──
    save_path = get_save_path(ppt_path)
    prs.save(save_path)
    print(f"\n{total_highlights} highlight(s) drawn.")
    print(f"Saved: {save_path}")

    # ── Open automatically ───────────────────────────────
    open_file(save_path)
    return save_path


# ─────────────────────────────────────────────
# STEP 6: Speech-triggered highlight
# ─────────────────────────────────────────────

def highlight_charts_by_speech(spoken_text, ppt_path):
    """
    Maps spoken words to finding types and highlights accordingly.

    "show the highest point"      → peak
    "where is the biggest drop"   → trend_down
    "revenue is going up"         → trend_up
    "which is the largest slice"  → biggest_slice
    "show the lowest value"       → trough
    anything else                 → all patterns
    """
    spoken        = spoken_text.lower()
    finding_types = []

    if any(w in spoken for w in [
        "highest", "peak", "maximum", "max", "top", "most"
    ]):
        finding_types.append("peak")

    if any(w in spoken for w in [
        "lowest", "minimum", "min", "bottom", "least", "trough"
    ]):
        finding_types.append("trough")

    if any(w in spoken for w in [
        "increase", "going up", "rise", "rising",
        "grew", "growth", "up", "higher", "trend up"
    ]):
        finding_types.append("trend_up")

    if any(w in spoken for w in [
        "decrease", "going down", "fall", "falling",
        "dropped", "decline", "down", "lower", "trend down"
    ]):
        finding_types.append("trend_down")

    if any(w in spoken for w in [
        "biggest", "largest", "majority", "dominant",
        "pie", "slice"
    ]):
        finding_types.append("biggest_slice")

    if not finding_types:
        print("No specific pattern detected — highlighting all.")
        finding_types = [
            "peak", "trough",
            "trend_up", "trend_down",
            "biggest_slice"
        ]

    print(f"Speech: '{spoken_text}' → highlighting: {finding_types}")
    return highlight_charts_in_ppt(ppt_path, finding_types)

