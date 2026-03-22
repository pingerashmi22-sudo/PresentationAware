from slides.chart_highlighter import highlight_charts_by_speech
from pptx import Presentation
from pptx.util import Inches
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE
import os

PPT_PATH = r"C:\Users\Heli P Shah\OneDrive\Documents\LOCAL SELF-GOVERNMENT - URBAN.pptx"

if not os.path.exists(PPT_PATH):
    print(f"File not found: {PPT_PATH}")
    exit()

# Check if PPT has real chart objects
prs         = Presentation(PPT_PATH)
chart_count = 0
for i, slide in enumerate(prs.slides):
    for shape in slide.shapes:
        if shape.has_chart:
            chart_count += 1
            print(f"Slide {i+1}: Chart found — '{shape.name}'")

# No charts found — create a test PPT with 3 chart types
if chart_count == 0:
    print("\nNo charts in this PPT.")
    print("Creating test_charts.pptx with bar, line and pie charts...\n")

    test_prs = Presentation()
    blank    = test_prs.slide_layouts[5]

    # Slide 1 — Bar chart
    s1  = test_prs.slides.add_slide(blank)
    cd1 = ChartData()
    cd1.categories = ["Q1", "Q2", "Q3", "Q4", "Q5"]
    cd1.add_series("Revenue",  (15, 45, 30, 70, 50))
    cd1.add_series("Expenses", (10, 20, 35, 25, 40))
    s1.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(1), Inches(1), Inches(8), Inches(5), cd1
    )
    print("Slide 1: Bar chart (Revenue vs Expenses)")

    # Slide 2 — Line chart
    s2  = test_prs.slides.add_slide(blank)
    cd2 = ChartData()
    cd2.categories = ["Jan", "Feb", "Mar", "Apr", "May", "Jun"]
    cd2.add_series("Growth", (10, 15, 12, 40, 35, 55))
    s2.shapes.add_chart(
        XL_CHART_TYPE.LINE,
        Inches(1), Inches(1), Inches(8), Inches(5), cd2
    )
    print("Slide 2: Line chart (Growth Jan-Jun)")

    # Slide 3 — Pie chart
    s3  = test_prs.slides.add_slide(blank)
    cd3 = ChartData()
    cd3.categories = ["Urban", "Rural", "Semi-Urban"]
    cd3.add_series("Population", (55, 30, 15))
    s3.shapes.add_chart(
        XL_CHART_TYPE.PIE,
        Inches(1), Inches(1), Inches(6), Inches(5), cd3
    )
    print("Slide 3: Pie chart (Urban/Rural/Semi-Urban)\n")

    test_path = os.path.join(
        os.path.dirname(PPT_PATH), "test_charts.pptx"
    )
    test_prs.save(test_path)
    print(f"Test PPT saved: {test_path}\n")
    PPT_PATH = test_path

# Live test loop
print("=== Type a phrase to highlight charts ===")
print("Try: 'show the highest point'")
print("Try: 'where is the biggest drop'")
print("Try: 'revenue is going up'")
print("Try: 'which is the largest slice'")
print("Try: 'show everything'")
print("Type 'quit' to stop\n")

while True:
    phrase = input("Speak: ").strip()
    if phrase.lower() == "quit":
        break
    if phrase:
        result = highlight_charts_by_speech(phrase, PPT_PATH)
        if result:
            print(f"PowerPoint should open automatically now.\n")
        else:
            print("No highlights drawn.\n")

