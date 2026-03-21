import os
import subprocess
import platform
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE


# ─────────────────────────────────────────────
# HELPERS
# ─────────────────────────────────────────────

def px_to_emu(px):
    return int(px * 914400 / 96)

def emu_to_px(emu):
    return round(emu / 914400 * 96)

def open_file(path):
    try:
        if platform.system() == "Windows":
            os.startfile(path)
        elif platform.system() == "Darwin":
            subprocess.call(["open", path])
        else:
            subprocess.call(["xdg-open", path])
        print(f"Opened: {path}")
    except Exception as e:
        print(f"Could not auto-open: {e}")
        print(f"Manually open: {path}")

def get_save_path(ppt_path):
    project_dir = os.path.join(
        os.path.expanduser("~"),
        "OneDrive", "Desktop", "slide"
    )
    os.makedirs(project_dir, exist_ok=True)
    filename  = os.path.basename(ppt_path).replace(
        ".pptx", "_highlighted.pptx"
    )
    return os.path.join(project_dir, filename)


# ─────────────────────────────────────────────
# STEP 1: Detect chart type
# ─────────────────────────────────────────────

def detect_chart_type(chart):
    try:
        ct = chart.chart_type
        type_map = {
            XL_CHART_TYPE.BAR_CLUSTERED:    "bar",
            XL_CHART_TYPE.BAR_STACKED:       "bar",
            XL_CHART_TYPE.COLUMN_CLUSTERED: "bar",
            XL_CHART_TYPE.COLUMN_STACKED:   "bar",
            XL_CHART_TYPE.LINE:             "line",
            XL_CHART_TYPE.LINE_MARKERS:     "line",
            XL_CHART_TYPE.LINE_STACKED:     "line",
            XL_CHART_TYPE.PIE:              "pie",
            XL_CHART_TYPE.PIE_EXPLODED:     "pie",
            XL_CHART_TYPE.DOUGHNUT:         "pie",
            XL_CHART_TYPE.AREA:             "area",
            XL_CHART_TYPE.XY_SCATTER:       "scatter",
        }
        return type_map.get(ct, "other")
    except Exception:
        return "other"


# ─────────────────────────────────────────────
# STEP 2: Extract chart data
# ─────────────────────────────────────────────

def extract_charts_from_slide(slide, slide_number):
    charts = []
    for shape in slide.shapes:
        if not shape.has_chart:
            continue

        chart      = shape.chart
        chart_type = detect_chart_type(chart)

        x = emu_to_px(shape.left)   if shape.left   else 0
        y = emu_to_px(shape.top)    if shape.top    else 0
        w = emu_to_px(shape.width)  if shape.width  else 0
        h = emu_to_px(shape.height) if shape.height else 0

        series_data = []
        try:
            for series in chart.series:
                try:
                    values = [float(v) if v is not None else 0.0 for v in series.values]
                except Exception:
                    values = []

                try:
                    cats = chart.plots[0].series[0].data_labels
                    categories = [str(c) for c in cats]
                except Exception:
                    categories = [f"Point {i+1}" for i in range(len(values))]

                series_data.append({
                    "name": str(series.name) if series.name else "Series",
                    "values": values,
                    "categories": categories
                })
        except Exception as e:
            print(f"  Could not read series on slide {slide_number}: {e}")

        charts.append({
            "slide_number": slide_number,
            "shape_name": shape.name,
            "chart_type": chart_type,
            "x": x, "y": y, "w": w, "h": h,
            "series": series_data
        })
    return charts


# ─────────────────────────────────────────────
# STEP 3: Analyze chart patterns
# ─────────────────────────────────────────────

def analyse_chart(chart_info):
    findings = []
    chart_type = chart_info.get("chart_type", "other")

    for series in chart_info.get("series", []):
        values = series.get("values", [])
        categories = series.get("categories", [])
        name = series.get("name", "Series")

        if not values or len(values) < 2:
            continue

        # Peak detection
        max_val = max(values)
        max_idx = values.index(max_val)
        findings.append({
            "type": "peak",
            "description": f"Highest value: {max_val} at {categories[max_idx]}"
        })

        # Trough detection
        min_val = min(values)
        min_idx = values.index(min_val)
        findings.append({
            "type": "trough",
            "description": f"Lowest value: {min_val} at {categories[min_idx]}"
        })

        # Pie slice detection
        if chart_type == "pie":
            biggest = max(values)
            biggest_idx = values.index(biggest)
            findings.append({
                "type": "biggest_slice",
                "description": f"Biggest slice: {categories[biggest_idx]}"
            })

    return findings


# ─────────────────────────────────────────────
# NEW STEP: Text-based Highlighting (Member 4 Task)
# ─────────────────────────────────────────────

def highlight_text_on_slide(slide, target_word):
    """
    MEMBER 4 TASK: Find a specific word on a slide and highlight it[cite: 41].
    Updates visual_highlighter.py to receive a word from the LLM[cite: 41].
    """
    yellow_highlight = RGBColor(255, 255, 0)
    found = False
    
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
            
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                if target_word.lower() in run.text.lower():
                    # Member 4: Apply visual change to font background [cite: 41]
                    run.font.highlight_color = yellow_highlight
                    print(f"✨ [Member 4] Highlighted '{target_word}' in: {shape.name}")
                    found = True
    return found


# ─────────────────────────────────────────────
# STEP 4: Draw highlight on Chart
# ─────────────────────────────────────────────

def draw_highlight_on_chart(slide, chart_info, finding):
    color_map = {
        "peak": RGBColor(0x00, 0xCC, 0x00),
        "trough": RGBColor(0xFF, 0x33, 0x33),
        "biggest_slice": RGBColor(0xFF, 0xCC, 0x00),
    }

    color = color_map.get(finding.get("type"), RGBColor(0xFF, 0x00, 0x00))
    cx, cy = px_to_emu(chart_info["x"]), px_to_emu(chart_info["y"])
    cw, ch = px_to_emu(chart_info["w"]), px_to_emu(chart_info["h"])

    border = slide.shapes.add_shape(1, cx, cy, cw, ch)
    border.fill.background()
    border.line.color.rgb = color
    border.line.width = Pt(4)

    label_box = slide.shapes.add_textbox(cx, cy - px_to_emu(40), cw, px_to_emu(36))
    tf = label_box.text_frame
    p = tf.paragraphs[0]
    p.text = finding.get("description", "")
    p.runs[0].font.size = Pt(12)
    p.runs[0].font.bold = True
    p.runs[0].font.color.rgb = color


# ─────────────────────────────────────────────
# STEP 5: Speech-triggered highlight (Master)
# ─────────────────────────────────────────────

def highlight_by_speech(spoken_text, ppt_path, current_slide_index=0):
    """
    Member 4 logic: Receive words from AI and find them on screen[cite: 41].
    """
    prs = Presentation(ppt_path)
    slide = prs.slides[current_slide_index]
    spoken = spoken_text.lower()
    
    # 1. Check for Chart patterns first
    charts = extract_charts_from_slide(slide, current_slide_index + 1)
    chart_highlighted = False
    
    if charts:
        for chart_info in charts:
            findings = analyse_chart(chart_info)
            for f in findings:
                if f["type"] in spoken or "show" in spoken:
                    draw_highlight_on_chart(slide, chart_info, f)
                    chart_highlighted = True

    # 2. If no chart action was taken, try text highlighting [cite: 41]
    if not chart_highlighted:
        # Extract the likely subject or use the whole phrase
        highlight_text_on_slide(slide, spoken_text)

    save_path = get_save_path(ppt_path)
    prs.save(save_path)
    open_file(save_path)
    return save_path