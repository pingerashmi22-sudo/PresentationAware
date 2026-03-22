from pptx import Presentation


def load_ppt(file_path):
    prs = Presentation(file_path)

    slides_data = []

    for i, slide in enumerate(prs.slides):
        slide_info = {
            "slide_number": i + 1,
            "title": "",
            "content": [],
            "notes": "",
            "full_text": ""
        }

        # Extract title
        if slide.shapes.title and slide.shapes.title.text:
            slide_info["title"] = slide.shapes.title.text.strip()

        # Extract bullet points / content
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()

                    # Avoid duplicates, empty text, and title repetition
                    if (
                        text
                        and text != slide_info["title"]
                        and text not in slide_info["content"]
                    ):
                        slide_info["content"].append(text)

        # Extract notes
        if slide.has_notes_slide:
            notes_text = slide.notes_slide.notes_text_frame.text.strip()
            slide_info["notes"] = notes_text

        # Create full_text (AI-friendly combined context)
        slide_info["full_text"] = " ".join(
            [slide_info["title"]] +
            slide_info["content"] +
            ([slide_info["notes"]] if slide_info["notes"] else [])
        )

        slides_data.append(slide_info)

    return slides_data